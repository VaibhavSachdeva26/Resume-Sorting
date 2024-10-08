from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import os

def extract_text(file):
    if file.type == "application/pdf":
        return extract_pdf_text(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_docx_text(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_ppt_text(file)


def extract_pdf_text(file):
    print("inside extract pdf text method..")

    reader = PdfReader(file)
    pdf_text = ""

    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        pdf_text += page.extract_text() + "\n"

        return pdf_text


def extract_docx_text(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_ppt_text(file):
    # Load the presentation
    prs = Presentation(file)

    # Initialize a list to store the text
    text_runs = []

    # Iterate through each slide in the presentation
    for slide in prs.slides:
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    ppt_text = ''
    for text in text_runs:
        ppt_text = ppt_text + os.linesep + text
    return ppt_text
